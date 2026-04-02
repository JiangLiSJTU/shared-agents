"""
insert_diagrams.py
CLI tool to insert images (PNG/JPEG) into specific slides of a PPTX file.
Part of the pptx-diagram-inserter skill.
"""

import argparse
import json
import os
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

def fit_dims(img_path, max_w, max_h):
    with Image.open(img_path) as img:
        nat_w, nat_h = img.size
    ratio = nat_w / nat_h
    if max_w / ratio <= max_h:
        w = max_w
        h = Emu(int(w / ratio))
    else:
        h = max_h
        w = Emu(int(h * ratio))
    return w, h

def main():
    parser = argparse.ArgumentParser(description="Insert images into PPTX slides.")
    parser.add_argument("--in-pptx", required=True, help="Input PPTX file")
    parser.add_argument("--out-pptx", required=True, help="Output PPTX file")
    parser.add_argument("--map", required=True, help="JSON string or path to JSON file mapping 'slide_index' (1-based) to 'image_path'")
    parser.add_argument("--top-with-img", type=float, default=3.8, help="Top margin (inches) if slide has existing image (default 3.8)")
    parser.add_argument("--top-no-img", type=float, default=2.3, help="Top margin (inches) if slide has no image (default 2.3)")
    parser.add_argument("--max-width", type=float, default=14.8, help="Max width in inches (default 14.8)")
    parser.add_argument("--max-height", type=float, default=5.0, help="Max height in inches (default 5.0)")
    parser.add_argument("--slide-width", type=float, default=16.0, help="Slide width in inches (default 16.0)")
    parser.add_argument("--slide-height", type=float, default=9.0, help="Slide height in inches (default 9.0)")
    args = parser.parse_args()

    if os.path.isfile(args.map):
        with open(args.map, "r", encoding="utf-8") as f:
            slide_map = json.load(f)
    else:
        slide_map = json.loads(args.map)

    # Convert mapping keys to int
    slide_map = {int(k): v for k, v in slide_map.items()}

    prs = Presentation(args.in_pptx)
    slides = prs.slides

    SLIDE_W = Inches(args.slide_width)
    SLIDE_H = Inches(args.slide_height)
    DIAG_MAX_W = Inches(args.max_width)
    DIAG_MAX_H = Inches(args.max_height)
    TOP_WITH_IMG = Inches(args.top_with_img)
    TOP_NO_IMG = Inches(args.top_no_img)

    inserted = 0
    skipped = 0

    for slide_num, img_path in sorted(slide_map.items()):
        if not os.path.exists(img_path):
            print(f"  [SKIP] Slide {slide_num:02d}: Not found -> {img_path}")
            skipped += 1
            continue

        if slide_num < 1 or slide_num > len(slides):
            print(f"  [SKIP] Slide {slide_num:02d}: Out of bounds.")
            skipped += 1
            continue

        slide = slides[slide_num - 1]

        pic_count = sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)

        w, h = fit_dims(img_path, DIAG_MAX_W, DIAG_MAX_H)
        top = TOP_WITH_IMG if pic_count > 0 else TOP_NO_IMG

        # Clamp height to avoid overflowing the slide boundaries
        max_h = SLIDE_H - Inches(0.15) - top
        if h > max_h:
            ratio = w.emu / h.emu
            h = max_h
            w = Emu(int(h.emu * ratio))

        left = (SLIDE_W - w) // 2

        try:
            slide.shapes.add_picture(img_path, left, top, w, h)
            img_name = os.path.basename(img_path)
            print(f"  [OK]   Slide {slide_num:02d}: {img_name} ({w.inches:.1f}\" x {h.inches:.1f}\")")
            inserted += 1
        except Exception as e:
            print(f"  [ERR]  Slide {slide_num:02d}: {os.path.basename(img_path)} -> {e}")
            skipped += 1

    prs.save(args.out_pptx)
    print("=" * 60)
    print(f"Done: {inserted} inserted / {skipped} skipped")
    print(f"Saved: {args.out_pptx}")

if __name__ == "__main__":
    main()
